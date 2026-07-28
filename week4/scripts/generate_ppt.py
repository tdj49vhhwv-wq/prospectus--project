#!/usr/bin/env python3
"""
生成 Week4 Presentation PPT: 三协电机 PE/VC 深讲
主题: 定位+抽取 两步独立方法论
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from config import *

# ── 颜色方案 ──
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
YELLOW_ACCENT = RGBColor(0xF3, 0x9C, 0x12)


def add_slide_number(slide, num, total=8):
    """添加页码"""
    left = Inches(4.5)
    top = Inches(7.1)
    txBox = slide.shapes.add_textbox(left, top, Inches(1.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT


def add_dark_bg(slide):
    """深色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE


def add_light_bg(slide):
    """浅色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_title_bar(slide, title_text, subtitle_text=None):
    """顶部标题栏"""
    # 深色标题条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # 标题文字
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8.5), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.78), Inches(8.5), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)


def add_body_text(slide, texts, left=0.6, top=1.5, width=8.8, font_size=14, spacing=1.2):
    """正文bullet列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item["text"]
        p.font.size = Pt(item.get("size", font_size))
        p.font.color.rgb = item.get("color", DARK_GRAY)
        p.space_after = Pt(item.get("after", 6))

        if item.get("bold"):
            p.font.bold = True
        if item.get("indent", 0) > 0:
            p.level = item["indent"]
        if item.get("alignment"):
            p.alignment = item["alignment"]


def add_table(slide, headers, rows, left=0.6, top=2.2, width=8.8, row_height=0.38):
    """添加表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(row_height * n_rows))
    table = table_shape.table

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER

    # 数据行
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT


def add_code_block(slide, code_text, left=0.6, top=3.0, width=8.8, height=3.0):
    """代码块"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0xA0, 0xD0, 0x40)
    p.font.name = "Courier New"


# ============================================================
# Slide 1: 封面
# ============================================================
def build_slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_dark_bg(slide)

    # 公司名
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "三协电机（920100）"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(8), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "PE/VC 深讲：定位 + 抽取 两步独立方法论"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER

    # 分隔线
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.6), Inches(3), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # 学生信息
    info_lines = [
        ("赵秉清", 18, True),
        ("Week 4 Presentation  |  2026-07-10", 13, False),
        ("核心主题：定向发行只披露合计，是否拆分需要判断", 12, False),
        ("直接披露与反推值的区别", 12, False),
    ]
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(8), Inches(2.5))
    tf3 = txBox3.text_frame
    for i, (txt, sz, bld) in enumerate(info_lines):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = WHITE if bld else RGBColor(0xBB, 0xCC, 0xDD)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    add_slide_number(slide, 1)


# ============================================================
# Slide 2: 进展总结
# ============================================================
def build_slide_progress(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "进展总结：Week3 → Week4", "从全量投喂到定位+抽取两步独立")

    headers = ["维度", "Week 3（旧方法）", "Week 4（新方法）"]
    rows = [
        ["方法", "全量PDF投喂LLM", "代码定位 → 精准截取10页 → 规则/JSON抽取"],
        ["三协电机覆盖", "5条认缴 + 1条快照", "完整9章节原文 + PE基金详情（备案/GP/LP）"],
        ["成本/次", "~¥3（382页全量）", "~¥0.02（150倍节省）"],
        ["定位方式", "PDF全文关键词扫描", "章节标题 → 段落PE/VC关键词 → 合并重叠"],
        ["输出格式", "正则硬编码提取", "JSON配置驱动（换公司只改配置）"],
        ["新增", "—", "PE基金结构化 / 15人拆分 / 募资差异检测 / 转增解析"],
    ]
    add_table(slide, headers, rows, left=0.4, top=1.5, width=9.2, row_height=0.45)

    # 关键总结
    add_body_text(slide, [
        {"text": "关键转变：从\"黑箱操作\"到\"工程化设计\"（闫老师要求）", "size": 13, "color": ACCENT_BLUE, "bold": True},
        {"text": "• 定位和抽取两步独立：先代码定位章节，再精准截取原文，最后JSON提取", "size": 12, "indent": 0},
        {"text": "• 成本从¥3/次降到¥0.02/次，且确定性100%（无LLM幻觉风险）", "size": 12, "indent": 0},
    ], top=5.8)

    add_slide_number(slide, 2)


# ============================================================
# Slide 3: 方法论核心 — 定位+抽取
# ============================================================
def build_slide_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "方法论核心：定位 + 抽取 两步独立", "代码定位章节 → 精准截取原文 → JSON配置抽取")

    # 流程图
    flow_items = [
        ("📄 PDF\n382页", 0.5, DARK_BLUE),
        ("🔍 定位\n代码关键词匹配\n5章节 10页", 3.2, ACCENT_BLUE),
        ("✂️ 截取\nPE/VC原文\n~20K字符", 5.9, ACCENT_BLUE),
        ("📊 抽取\nJSON配置驱动\n6条结构化记录", 8.6, GREEN_ACCENT),
    ]

    for text, left, color in flow_items:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(2.2), Inches(1.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # 箭头
    for x in [2.8, 5.5, 8.2]:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(2.3), Inches(0.35), Inches(0.3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_GRAY
        shape.line.fill.background()

    # 两列对比
    # 左列: 定位
    add_body_text(slide, [
        {"text": "Step 1: 定位（代码化）", "size": 16, "color": ACCENT_BLUE, "bold": True},
        {"text": "", "size": 6},
        {"text": "章节级关键词：", "size": 12, "bold": True},
        {"text": "\"发行人基本情况\" / \"发行融资\"", "size": 11, "indent": 1},
        {"text": "\"股权结构\" / \"股东及实际控制人\"", "size": 11, "indent": 1},
        {"text": "\"私募投资基金纳入监管\"", "size": 11, "indent": 1},
        {"text": "段落级关键词：", "size": 12, "bold": True},
        {"text": "\"定向发行\" / \"增资\" / \"备案编码\"", "size": 11, "indent": 1},
        {"text": "→ 从382页精准定位到10页", "size": 12, "color": GREEN_ACCENT, "bold": True},
    ], left=0.5, top=3.5, width=4.2)

    # 右列: 抽取
    add_body_text(slide, [
        {"text": "Step 2: 抽取（JSON配置化）", "size": 16, "color": ACCENT_BLUE, "bold": True},
        {"text": "", "size": 6},
        {"text": "不用正则，用JSON配置驱动：", "size": 12, "bold": True},
        {"text": "find_between(\"已收到\", \"缴纳的出资款\")", "size": 11, "indent": 1},
        {"text": "find_number_after(\"发行价格为\", \"元/股\")", "size": 11, "indent": 1},
        {"text": "smart_split_names(名单字符串)", "size": 11, "indent": 1},
        {"text": "→ 6条结构化记录，换公司只改JSON", "size": 12, "color": GREEN_ACCENT, "bold": True},
    ], left=5.2, top=3.5, width=4.2)

    add_slide_number(slide, 3)


# ============================================================
# Slide 4: 三协电机时间线
# ============================================================
def build_slide_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "三协电机（920100）股本变化时间线", "从有限公司设立到IPO前，两次定增 + 资本公积转增")

    events = [
        ("2002.11", "有限公司设立\n盛祎+朱绶青", "⚠ t0出资额\n未在招股书披露", YELLOW_ACCENT),
        ("2022.02", "新三板挂牌\n基础层", "证券代码873669", DARK_GRAY),
        ("2022.08", "★ 第一次定增\n稳正景明+长泽创投", "4.48元/股\n530万股 2374万", ACCENT_BLUE),
        ("2023.05", "调入创新层", "", DARK_GRAY),
        ("2023.09", "★ 第二次定增\n盛祎等15名自然人", "5.41元/股\n321.5万股 1723万", ACCENT_BLUE),
        ("2023.12", "10转增3.8股\n转增1462万股", "总股本→5310.93万", GREEN_ACCENT),
        ("2025.07", "IPO前\nPE持股: 稳正景明9.16%\n长泽创投4.61%", "招股书签署日", DARK_BLUE),
    ]

    for i, (date, title, detail, color) in enumerate(events):
        y = 1.5 + i * 0.78

        # 圆点
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.15), Inches(0.2), Inches(0.2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        # 连线
        if i < len(events) - 1:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.88), Inches(y + 0.35), Inches(0.04), Inches(0.45)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            line.line.fill.background()

        # 日期
        txBox = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(1.2), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = date
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        # 标题
        txBox2 = slide.shapes.add_textbox(Inches(2.6), Inches(y), Inches(3.5), Inches(0.55))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = DARK_GRAY

        # 详情
        if detail:
            txBox3 = slide.shapes.add_textbox(Inches(6.2), Inches(y), Inches(3.3), Inches(0.55))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = detail
            p3.font.size = Pt(10)
            p3.font.color.rgb = DARK_GRAY

    add_slide_number(slide, 4)


# ============================================================
# Slide 5: PDF原文证据
# ============================================================
def build_slide_evidence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "PDF 原文证据", "关键页码：p31-32 定增 / p35 股东表 / p39 私募基金备案")

    # 证据 1
    add_body_text(slide, [
        {"text": "证据1: 2022年定增 — PE进入 (p31-32)", "size": 13, "color": ACCENT_BLUE, "bold": True},
        {"text": "\"本次股票发行价格为4.48元/股，共发行普通股530.00万股，募集资金总额为2,374.40万元。", "size": 10, "indent": 0},
        {"text": "经审验，截至2022年8月9日，公司已收到稳正景明、长泽创投缴纳的出资款2,374.40万元。\"", "size": 10, "indent": 0},
        {"text": "⚠ 关键问题：2374.40万是两家合计，PDF未单独披露各自出资金额", "size": 10, "color": RED_ACCENT, "bold": True, "indent": 0},
        {"text": "", "size": 6},
        {"text": "证据2: 2023年定增 — 拟募集≠实际收到 (p32)", "size": 13, "color": ACCENT_BLUE, "bold": True},
        {"text": "\"本次股票拟发行价格为5.41元/股，拟发行普通股321.50万股，拟募集资金总额为1,739.32万元。", "size": 10, "indent": 0},
        {"text": "已收到盛祎、盛松...等15名认购人缴纳的出资款1,723.09万元。\"", "size": 10, "indent": 0},
        {"text": "⚠ price×shares=1739万 ≠ 实际1723万（差16.23万/0.93%）→ 待复核", "size": 10, "color": RED_ACCENT, "bold": True, "indent": 0},
    ], top=1.5)

    # 证据 3: PE备案
    add_body_text(slide, [
        {"text": "证据3: 私募基金备案 (p39)", "size": 13, "color": ACCENT_BLUE, "bold": True},
        {"text": "\"稳正景明于2020年11月16日在中国证券投资基金业协会备案，备案编码为SNG030；", "size": 10, "indent": 0},
        {"text": "长泽创投于2022年6月27日备案，备案编码为SVU935。", "size": 10, "indent": 0},
        {"text": "基金管理人为深圳市稳正资产管理有限公司，备案登记编号P1003586。\"", "size": 10, "indent": 0},
    ], top=4.3)

    add_slide_number(slide, 5)


# ============================================================
# Slide 6: Gold表示例
# ============================================================
def build_slide_gold(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "Gold 表示例：结构化提取结果", "JSONL格式 — 直接披露 vs 反推值")

    # subscription_flow 示例
    add_body_text(slide, [
        {"text": "SubscriptionFlow（认缴流量）— 2022年定增", "size": 12, "color": ACCENT_BLUE, "bold": True},
    ], top=1.4)

    code_sf = (
        '{\n'
        '  "record_type": "subscription_flow",\n'
        '  "subscriber_name": "深圳市稳正景明创业投资企业(有限合伙)",\n'
        '  "subscription_date": "2022-08-09",\n'
        '  "price_per_share": 4.48,\n'
        '  "shares_subscribed": null,           ← ⚠ 未披露\n'
        '  "amount_subscribed": null,            ← ⚠ 仅合计\n'
        '  "event_context": "增资",\n'
        '  "notes": "PDF只披露稳正景明+长泽创投合计2374.40万元，\n'
        '            未单独披露各自金额，不能反推"\n'
        '}'
    )
    add_code_block(slide, code_sf, top=1.8, width=8.8, height=2.2)

    # 对比表
    add_body_text(slide, [
        {"text": "直接披露 vs 反推 → 核心原则", "size": 13, "color": ACCENT_BLUE, "bold": True},
    ], top=4.2)
    headers = ["字段", "来源", "类型", "原则"]
    rows = [
        ["稳正景明持股 486.70万股 / 9.16%", "p37 前十名股东表", "✅ 直接披露", "照抄"],
        ["2022年发行价 4.48元/股", "p31", "✅ 直接披露", "照抄"],
        ["稳正景明单独出资金额", "—", "❌ 未披露", "留空 null"],
        ["长泽创投单独出资金额", "—", "❌ 未披露", "留空 null"],
        ["设立时盛祎/朱绶青出资额", "—", "❌ 未披露", "需查公开转让说明书"],
    ]
    add_table(slide, headers, rows, top=4.5, row_height=0.36)

    add_slide_number(slide, 6)


# ============================================================
# Slide 7: 与其他公司共性
# ============================================================
def build_slide_commonality(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "与其他公司的共性：从个案提炼方法", "三协电机反映了 IPO 招股书抽取中的共性问题")

    headers = ["共性", "三协电机的表现", "其他公司表现"]
    rows = [
        ["PE/VC 通过定向发行进入", "2022年稳正景明+长泽创投定增", "友升A/B轮增资、云汉芯城多轮增资"],
        ["私募基金备案信息在招股书", "p38-39 备案编码SNG030/SVU935", "科创板/创业板标配章节"],
        ["发行前股东结构表（核心来源）", "p35-36 完整19名股东 + IPO前后对比", "所有招股书第四节必有"],
        ["历史沿革不在招股书正文", "需查《公开转让说明书》", "影石创新VIE细节在问询回复"],
        ["单位一致性", "全招股书统一用\"万元/万股\"", "黄山谷捷有78/780十进位混淆"],
        ["股权代持", "2021年3月解除（p39）", "星图测控也有代持还原"],
        ["定向发行只给合计", "2374万=稳正景明+长泽创投(未拆分)", "云汉芯城英文投资人也有类似问题"],
    ]
    add_table(slide, headers, rows, left=0.4, top=1.5, width=9.2, row_height=0.48)

    # 可迁移方法
    add_body_text(slide, [
        {"text": "可迁移方法：换一家类似公司，应该先去招股书哪里找信息？", "size": 14, "color": ACCENT_BLUE, "bold": True},
        {"text": "北交所公司：第四节\"发行人基本情况\" → \"报告期内发行融资\" + \"股东及实际控制人\" + \"股本情况\"", "size": 12, "indent": 0},
        {"text": "科创板/创业板：第五节\"发行人基本情况\" → \"股本演变\" + \"历次增资\" + \"新增股东\"", "size": 12, "indent": 0},
        {"text": "共性原则：数字能直接抄就抄，PDF没写的留空不反推，差异标\"待复核\"不硬凑", "size": 12, "color": RED_ACCENT, "bold": True, "indent": 0},
    ], top=5.5)

    add_slide_number(slide, 8)


# ============================================================
# Slide 8: 自动化流程与失败点
# ============================================================
def build_slide_automation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_title_bar(slide, "自动化流程与失败点分析", "规则能做什么 / LLM该做什么 / 人工必须做什么")

    # 流程
    add_body_text(slide, [
        {"text": "当前流程：PDF(382页) → 代码定位 → 10页原文 → JSON配置抽取 → JSONL", "size": 12, "color": DARK_BLUE, "bold": True},
    ], top=1.4)

    headers = ["环节", "成功?", "方法", "说明"]
    rows = [
        ["PDF→文本", "✅", "PyMuPDF", "文本层完整"],
        ["章节定位", "✅", "keyword match", "\"发行融资\"等标题明确"],
        ["PE投资人识别", "✅", "string find", "稳正景明/长泽创投 → 含\"创业投资\""],
        ["认购金额提取", "⚠", "find_between", "合计无法拆分 → 留空，不能反推"],
        ["15人名单拆分", "✅", "smart_split", "分隔符+和 → 15/15 验证通过"],
        ["PE基金备案", "✅", "fixed_value", "SNG030/SVU935/P1003586"],
        ["表格提取", "⚠", "PyMuPDF tables", "合并单元格 → 列映射偏差"],
        ["设立事件", "❌", "—", "招股书本身不披露 → 外部数据"],
    ]
    add_table(slide, headers, rows, top=2.0, row_height=0.32)

    # 三层分工
    add_body_text(slide, [
        {"text": "规则 → 定位章节 / 识别名称 / 数字提取     LLM → 语义消歧 / 差异判断     人工 → PDF验证 / 不反推", "size": 12, "color": ACCENT_BLUE, "bold": True},
    ], top=5.6)

    add_body_text(slide, [
        {"text": "核心教训：PDF没写的就留空，标'待复核'，不硬凑。这就是工程化的底线。", "size": 11, "color": RED_ACCENT, "bold": True},
    ], top=6.2)

    add_slide_number(slide, 7)


# ============================================================
# Slide 9: 下一步计划
# ============================================================
def build_slide_next(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(slide)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "下一步计划"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    headers = ["优先级", "任务", "说明"]
    rows = [
        ["P0", "补充公开转让说明书", "获取设立时出资结构（t0）"],
        ["P0", "15人→15条SF记录", "每人一行subscription_flow"],
        ["P0", "PE单独认购金额", "标注\"PDF未单独披露\"，不反推"],
        ["P1", "跨页表格精确解析", "改进PyMuPDF合并单元格映射"],
        ["P1", "Cross-check增强", "总股本×持股比例≈持股数校验"],
        ["P2", "8家公司比对验证", "北交所 vs 科创板招股书结构差异"],
    ]
    add_table(slide, headers, rows, left=1.0, top=1.8, width=8.0, row_height=0.52)

    # 底部总结
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.2))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "一句话总结："
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf2.add_paragraph()
    p3.text = "招股书定向发行常只披露认购合计，不要为了填满表格去反推。\nPDF没写的就留空，标\"待复核\"，这才是可回源、可信赖的工程化方法。"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(0xF3, 0x9C, 0x12)
    p3.alignment = PP_ALIGN.CENTER

    add_slide_number(slide, 8)


# ============================================================
# 主函数
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    build_slide_cover(prs)
    build_slide_progress(prs)
    build_slide_methodology(prs)
    build_slide_timeline(prs)
    build_slide_evidence(prs)
    build_slide_gold(prs)
    build_slide_commonality(prs)
    build_slide_automation(prs)
    build_slide_next(prs)

    output_path = OUTPUTS_DIR / "week4_三协电机_PEVC_深讲.pptx"
    prs.save(str(output_path))
    print(f"✓ PPT 已生成: {output_path}")
    print(f"  共 {len(prs.slides)} 页")
    return 0


if __name__ == "__main__":
    sys.exit(main())
